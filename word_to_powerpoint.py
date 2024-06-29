from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import os
from datetime import datetime


doc_path = '/path/to/your/Word/document.docx'
doc = Document(doc_path)


text_content = []
for para in doc.paragraphs:
    text_content.append(para.text)


prs = Presentation()
prs.slide_width = Inches(10)  
prs.slide_height = Inches(7.5)  


for para in text_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  

    
    content_shape = slide.placeholders[1]
    content_shape.text = para
    content_shape.text_frame.paragraphs[0].font.size = Pt(18) 


downloads_folder = os.path.join(os.getcwd(), 'Downloads')
timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
output_filename = f'output_{timestamp}.pptx'
output_path = os.path.join(downloads_folder, output_filename)


prs.save(output_path)

print(f"Presentation saved as '{output_path}'")


os.system(f'open "{output_path}"')
